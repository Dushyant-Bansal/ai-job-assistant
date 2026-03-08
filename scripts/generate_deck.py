"""Generate PowerPoint deck for AI Job Training Assistant design and testing."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml

# Formal modern color scheme: slate navy + teal accent
FONT_HEADING = "Segoe UI"
FONT_BODY = "Segoe UI"
FONT_MONO = "Consolas"

COLOR_NAVY = RGBColor(0x0F, 0x17, 0x2A)       # Primary headings
COLOR_SLATE = RGBColor(0x33, 0x44, 0x55)       # Body text
COLOR_MUTED = RGBColor(0x64, 0x74, 0x8B)      # Subtitle, secondary
COLOR_ACCENT = RGBColor(0x0D, 0x94, 0x88)     # Teal accent
COLOR_BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)   # Light background
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _add_accent_bar(slide, prs) -> None:
    """Add a thin vertical accent bar on the left edge."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.08), prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()


def _add_header_bar(slide) -> None:
    """Add a subtle header bar across the top."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.15)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_NAVY
    shape.line.fill.background()


def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide)
    _add_accent_bar(slide, prs)

    left = Inches(0.75)
    top = Inches(2.2)
    width = Inches(8.5)
    height = Inches(1.5)
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = FONT_HEADING
    p.font.color.rgb = COLOR_NAVY
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.name = FONT_BODY
        p2.font.color.rgb = COLOR_MUTED
        p2.space_before = Pt(14)


def add_content_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide)
    _add_accent_bar(slide, prs)

    left = Inches(0.75)
    top = Inches(0.6)
    width = Inches(8.5)
    height = Inches(0.9)
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = FONT_HEADING
    p.font.color.rgb = COLOR_NAVY

    left = Inches(0.75)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(5.5)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Ensure bullet prefix for list items
        text = bullet.strip()
        if text and not text.startswith("•") and not text.startswith("-"):
            text = "• " + text
        p.text = text
        p.font.size = Pt(16)
        p.font.name = FONT_BODY
        p.font.color.rgb = COLOR_SLATE
        p.space_before = Pt(6)
        p.level = 0


def _add_arrow_connector(slide, x1: float, y1: float, x2: float, y2: float) -> None:
    """Add a connector with arrowhead from (x1,y1) to (x2,y2) in inches.
    Uses tailEnd so arrow points from source toward target (headEnd points opposite in OOXML).
    """
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.fill.solid()
    conn.line.fill.fore_color.rgb = COLOR_SLATE
    conn.line.width = Pt(1.5)
    ln = conn.line._get_or_add_ln()
    ln.append(parse_xml(
        '<a:tailEnd type="arrow" w="med" len="med" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
    ))


def _add_node(slide, left: float, top: float, width: float, height: float, label: str, accent: bool = False) -> None:
    """Add a rounded rectangle node with centered text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT if accent else COLOR_BG_LIGHT
    shape.line.color.rgb = COLOR_ACCENT if accent else RGBColor(0xCB, 0xD5, 0xE1)
    shape.line.width = Pt(1) if accent else Pt(0.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.name = FONT_BODY
    tf.paragraphs[0].font.color.rgb = COLOR_NAVY if accent else COLOR_SLATE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_langgraph_diagram_slide(prs: Presentation, title: str) -> None:
    """Add architecture slide with shape-based LangGraph flow diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide)
    _add_accent_bar(slide, prs)

    left = Inches(0.75)
    top = Inches(0.5)
    width = Inches(8.5)
    height = Inches(0.7)
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.name = FONT_HEADING
    p.font.color.rgb = COLOR_NAVY

    # Layout: increased spacing between blocks
    w, h = 1.3, 0.35
    gap = 0.22
    row = 0.75  # vertical spacing between rows
    base_left = 0.9
    base_top = 1.35

    r0, r1, r2, r3, r4, r5, r6 = (
        base_top,
        base_top + row,
        base_top + 2 * row,
        base_top + 3 * row,
        base_top + 4 * row,
        base_top + 5 * row,
        base_top + 6 * row,
    )

    # Row 0: START
    _add_node(slide, base_left + 3.4, r0, 0.8, 0.3, "START", accent=True)
    # Row 1: parse_resume -> analyze_jd -> validate_industry
    _add_node(slide, base_left, r1, w, h, "parse_resume")
    _add_node(slide, base_left + w + gap, r1, w, h, "analyze_jd")
    _add_node(slide, base_left + 2 * (w + gap), r1, w, h, "validate_industry")
    # Row 2: normalize_skills
    _add_node(slide, base_left + 1.5, r2, w, h, "normalize_skills")
    # Row 3: match_skills
    _add_node(slide, base_left + 1.5, r3, w, h, "match_skills")
    # Row 4: parallel resources (6 boxes, spaced out)
    pw, ph = 1.0, 0.3
    res_gap = 0.2
    res_labels = ["web", "news", "yt", "amazon", "courses", "blog"]
    res_left = base_left  # 6 boxes: 6*pw + 5*res_gap = 6 + 1 = 7
    for i, lbl in enumerate(res_labels):
        _add_node(slide, res_left + i * (pw + res_gap), r4, pw, ph, lbl)
    # Row 5: generate_training_plan
    _add_node(slide, base_left + 1.5, r5, w, h, "generate_training_plan")
    # Row 6: END
    _add_node(slide, base_left + 3.4, r6, 0.8, 0.3, "END", accent=True)

    def cx(l, width): return l + width / 2
    def cy(t, height): return t + height / 2

    match_cx = cx(base_left + 1.5, w)
    match_bottom = r3 + h
    res_top = r4
    res_bottom = r4 + ph
    gen_top = r5
    gen_cx = cx(base_left + 1.5, w)

    # START -> parse_resume
    _add_arrow_connector(slide, cx(base_left + 3.4, 0.8), r0 + 0.3, cx(base_left, w), r1)
    # parse_resume -> analyze_jd -> validate_industry
    _add_arrow_connector(slide, base_left + w, cy(r1, h), base_left + w + gap, cy(r1, h))
    _add_arrow_connector(slide, base_left + 2 * w + gap, cy(r1, h), base_left + 2 * (w + gap), cy(r1, h))
    # validate_industry -> normalize_skills
    _add_arrow_connector(slide, cx(base_left + 2 * (w + gap), w), r1 + h, cx(base_left + 1.5, w), r2)
    # normalize_skills -> match_skills
    _add_arrow_connector(slide, cx(base_left + 1.5, w), r2 + h, match_cx, r3)
    # match_skills -> each of web, news, yt, amazon, courses, blog (6 separate arrows)
    for i in range(6):
        res_cx = res_left + i * (pw + res_gap) + pw / 2
        _add_arrow_connector(slide, match_cx, match_bottom, res_cx, res_top)
    # each resource -> generate_training_plan (6 arrows)
    for i in range(6):
        res_cx = res_left + i * (pw + res_gap) + pw / 2
        _add_arrow_connector(slide, res_cx, res_bottom, gen_cx, gen_top)
    # generate_training_plan -> END
    _add_arrow_connector(slide, gen_cx, r5 + h, cx(base_left + 3.4, 0.8), r6)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "AI Job Training Assistant",
        "Design, Architecture & Testing Overview",
    )

    # Slide 2: Problem Statement
    add_content_slide(
        prs,
        "Problem Statement",
        [
            "Job seekers struggle to identify skill gaps between their resume and job requirements.",
            "Manually comparing a resume to a job description is time-consuming and prone to oversight.",
            "Finding relevant learning resources for identified gaps is fragmented across many platforms.",
            "There is no unified tool that automates gap analysis and curates learning resources.",
            "Solution: An AI-powered assistant that analyzes resume vs. JD, identifies gaps, "
            "and suggests personalized training resources.",
        ],
    )

    # Slide 3: Design Overview
    add_content_slide(
        prs,
        "Design Overview",
        [
            "AI-powered tool that compares a resume to a software job description.",
            "Identifies skill gaps and suggests personalized learning resources.",
            "Key capabilities:",
            "Resume analysis — extracts technical & soft skills, rates 1–100",
            "Job description analysis — required skills, industry, domain",
            "Skill matching — weighted match %, gap detection",
            "Resource discovery — web, news, YouTube, books, courses",
            "Training plan — prioritized, actionable study plan",
        ],
    )

    # Slide 3: Design Principles
    add_content_slide(
        prs,
        "Design Principles",
        [
            "Modular orchestration — LangGraph state machine with discrete nodes",
            "Structured outputs — Pydantic models for resume, JD, training plan",
            "Canonical skill mapping — synonyms (LLMs, RAG, K8s) normalized for matching",
            "Graceful degradation — missing API keys disable features, not the app",
            "Privacy-first — resume & JD not stored on any server",
            "User overrides — add skills, adjust ratings, change domain; re-analyze without re-upload",
        ],
    )

    # Slide 4: Technology Stack
    add_content_slide(
        prs,
        "Technology Stack",
        [
            "Orchestration: LangGraph (state graph)",
            "LLM: OpenAI GPT-4o via LangChain",
            "UI: Streamlit",
            "Web search: Tavily (articles, news, blog, books, courses)",
            "Video search: YouTube Data API v3",
            "Parsing: PyMuPDF (PDF), python-docx (DOCX), plain TXT",
        ],
    )

    # Slide 6: Architecture Diagram (shape-based)
    add_langgraph_diagram_slide(prs, "Architecture — LangGraph Flow")

    # Slide 6: State & Data Flow
    add_content_slide(
        prs,
        "State & Data Flow",
        [
            "GraphState (TypedDict) flows through all nodes",
            "Inputs: resume_text, job_description, ignore_programming_languages",
            "Resume/JD nodes: user_skills, required_skills, job_title, software_domain",
            "Matching: skill_canonical_map, match_percentage, skill_gaps",
            "Resources: web_articles, news, youtube, amazon, courses, blog (merged via reducers)",
            "Output: training_plan, resource_search_warnings",
        ],
    )

    # Slide 7: Unit Tests
    add_content_slide(
        prs,
        "Unit Tests (28 tests)",
        [
            "test_parsers — parse_txt, parse_resume (TXT, PDF, unsupported)",
            "test_config — normalize_skill_name, is_programming_language, missing_api_keys",
            "test_skill_matcher — perfect match, gap detection, synonym matching",
            "test_resource_finder — _top_gap_skill_names fallbacks",
            "test_export — build_export_json structure",
            "Run: pytest tests/ -v",
        ],
    )

    # Slide 8: LLM-as-a-Judge Eval Tests
    add_content_slide(
        prs,
        "LLM-as-a-Judge Eval Tests (3 tests)",
        [
            "Uses GPT-4o-mini to judge analysis quality",
            "Test 1: Skill gaps plausible for software engineer JD?",
            "Test 2: Training plan prioritizes critical gaps appropriately?",
            "Test 3: Synonym matching (LLMs vs Large Language Models) correct?",
            "Requires OPENAI_API_KEY; skip with SKIP_EVALS=1",
        ],
    )

    # Slide 9: Iterative Manual Testing & Code Updates (combined)
    add_content_slide(
        prs,
        "Iterative Manual Testing & Code Updates",
        [
            "Learning resources not showing — fixed stream processing to handle parallel nodes",
            "LLMs vs Large Language Models not matching — normalized canonical values in skill_normalizer",
            "Chart not reflecting skill matches — added _skill_canonical() for chart lookups",
            "Tavily returning 0 — surface API errors in resource_search_warnings",
            "Token limit (429) — truncate JD to 15,000 chars",
            "Empty skill_gaps when ignoring PLs — added required_skills_for_resources fallback",
            "Add skill to resume — override UI with Add/Remove, re-analyze from normalize_skills",
            "Export as JSON — build_export_json, download button",
            "Privacy note — displayed in sidebar",
        ],
    )

    out_path = Path(__file__).resolve().parent.parent / "AI_Job_Assistant_Design_Deck.pptx"
    prs.save(str(out_path))
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
